#!/usr/bin/env python3
"""Generate VoiceFlow Builder PPT (Chinese version)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BG = RGBColor(0x1e, 0x29, 0x3b)
WHITE = RGBColor(0xff, 0xff, 0xff)
ACCENT = RGBColor(0x4f, 0x46, 0xe5)
LIGHT_TEXT = RGBColor(0x94, 0xa3, 0xb8)
SUCCESS = RGBColor(0x22, 0xc5, 0x5e)
WARNING = RGBColor(0xf5, 0x9e, 0x0b)
ERROR = RGBColor(0xef, 0x44, 0x44)

def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=16, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return txBox

def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

def create_slide(prs, title_text, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide)
    # Accent line
    add_accent_line(slide, 0.8, 1.0, 1.5)
    # Title
    add_text_box(slide, 0.8, 1.1, 8.5, 0.8, title_text, font_size=28, bold=True, color=WHITE)
    # Subtitle
    if subtitle:
        add_text_box(slide, 0.8, 1.8, 8.5, 0.5, subtitle, font_size=14, color=LIGHT_TEXT)
    # Bullets
    bullet_top = 2.5 if subtitle else 2.2
    add_bullet_list(slide, 1.0, bullet_top, 8.3, 4.5, bullets, font_size=16, color=WHITE)
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slides_data = [
        ("语流 VoiceFlow Builder", ["10分钟搭建专业语音接待员", "面向中小企业的可视化工作流构建器", "2026年5月"], "面向中小企业的可视化工作流构建器SaaS产品"),
        ("中小企业的语音接待困境", ["78%的中小企业仍在用人工接听所有来电", "现有方案要么太贵（呼叫中心），要么太简陋（录音电话）", "非技术用户无法使用开发者工具（Bland AI/Vapi需API能力）", "核心诉求：10分钟能用 + 业务复杂能扛"]),
        ("语流的三大核心能力", ["可视化工作流构建器：拖拽节点+连线，零代码搭建流程", "规则优先AI增强：核心流程零延迟零成本，AI在三个质变点精准赋能", "节点级运营洞察：热力图+漏斗定位流失节点，数据驱动迭代"]),
        ("谁需要语流？", ["小企业主(1-20人)：无技术团队，需快速搭建语音接待员", "中型企业团队管理者(20-200人)：需客服分流和自动化", "在线业务运营者：需网站嵌入语音通话，无电话号码需求"]),
        ("工作流如何运转？", ["来电触发 → 节点执行 → 条件路由 → 动作执行 → 通话结束", "每个节点：输入 → 处理逻辑 → 输出，确定性执行", "AI节点有降级方案：AI不可用时自动回退到规则流程"]),
        ("三步构建语音接待员", ["1. 拖拽节点：从节点库拖入画布，19种节点覆盖全场景", "2. 连线配置：连线表达数据流，双击节点配置参数", "3. 模拟发布：文字模拟验证 → 一键发布上线"]),
        ("19个节点，四层递进", ["核心层(8节点)：来电触发/问候/提问/收集/转接/发短信/记录/结束", "逻辑层(4节点)：条件分支/If-Else/营业时间判断/来电类型检测", "集成层(4节点)：CRM查询/日历预订/发送通知/Webhook", "AI增强层(3节点)：意图检测/知识库问答/通话摘要"]),
        ("核心层 — 通话流程骨架", ["来电触发器：唯一入口，自动提取来电信息", "问候/开场白：按时间段智能选话术，变量替换", "提问节点：支持语音+按键双输入，超时重试可配", "转接呼叫：暖/冷转接，失败路径自动降级", "每个节点有明确配置项，用户知道配什么"]),
        ("逻辑层 — 流程控制中枢", ["条件分支：多条件路由(≤4分支)，支持正则匹配", "If/Else：简单二元判断，降低配置复杂度", "营业时间判断：时区+节假日日历+特殊日期覆盖", "来电类型检测：老客户/VIP/黑名单自动识别"]),
        ("规则优先、AI增强策略", ["17个规则节点覆盖80%场景，零延迟零增量成本", "AI仅在3个质变点引入：意图检测(通话时长-35%)/知识库问答(人工转接-68%)/通话摘要(节省100%人工)", "情感检测和智能路由归入P2：规则近似替代代价可控(精度差~15%)", "AI节点必须有降级方案：意图未识别→按键菜单"]),
        ("三个AI质变点", ["AI意图检测：将按键导航替换为自然语言理解，1步直达vs4.2步，放弃率6%vs28%", "知识库问答(RAG)：一个知识库覆盖数百FAQ，拦截68%人工转接", "通话摘要：离线生成，零实时延迟，3分钟人工→30秒AI"]),
        ("成本与延迟的精确边界", ["纯规则引擎：$0增量成本，<10ms延迟", "+AI意图检测：+$0.02/min，+300ms，回报：通话时长-35%", "+知识库问答：+$0.03/min，+600ms，回报：人工转接-68%", "全AI Agent：+$0.08-0.12/min，+800ms，仅高复杂度场景值得"]),
        ("三模通信接入", ["平台托管号码：开箱即用，购买→绑定→上线，$1-3/月+0.01/min", "自带号码(SIP/移植)：保持现有号码，SIP Forwarding或号码移植", "WebRTC浏览器通话：无需电话号码，嵌入网站直接通话，$0.003-0.008/min", "全市场唯一同时提供三种模式的语音AI平台"]),
        ("两种模拟测试", ["文字模拟：画布内即时验证，秒级迭代，高亮当前执行节点", "真实通话模拟：临时号码语音测试，TTS+ASR真实验证", "被低估的杀手功能：文字模拟将设计-测试循环从分钟级降到秒级"]),
        ("安全迭代，恐惧归零", ["每次保存自动快照，保留最近50个版本", "版本Diff：新增绿/删除红/修改黄高亮对比", "一键回滚：修改出错30秒恢复", "草稿/发布分离：已发布版本不可直接修改"]),
        ("节点级洞察，精准归因", ["基础指标：通话量/接通率/平均时长/满意度", "节点热力图：绿>90%/黄70-90%/红<70%，定位流失节点", "转化漏斗：来电→意图检测→信息收集→转接→结束", "异常告警：连续失败/超时/路由失败实时通知"]),
        ("示例：标准客服分流", ["来电→营业时间判断→(营业)AI意图检测→咨询/预约/投诉分支", "非营业→知识库FAQ或留言", "预约分支→日历预订→发送确认短信→结束", "关键配置：意图检测3意图+0.7阈值+按键降级"]),
        ("示例：医疗诊所高级自动化", ["来电→来电类型检测→(老患者)CRM查询→AI意图检测→预约/咨询/紧急分支", "紧急→暖转接急诊，跳过所有流程", "预约→匹配主治医生日历→冲突推荐最近时段", "知识库覆盖FAQ(营业时间/医保/停车/科室)"]),
        ("产品定位与竞争力", ["差异化①：三模通信接入+双模模拟测试（全市场唯一）", "差异化②：规则优先AI增强的分层节点体系", "差异化③：节点级运营洞察看板（热力图+漏斗归因）", "vs Bland/Vapi/Retell：非技术用户10分钟上手", "vs Synthflow：节点体系透明+三模接入+节点级洞察"]),
        ("产品蓝图与路线图", ["MVP：核心层+逻辑层+画布+平台托管号码+文字模拟+基础看板", "V1.0：+AI增强层3节点+集成层+自带号码+真实通话模拟+版本管理+节点热力图", "V2.0：+WebRTC+情感检测+智能路由+多租户RBAC+API开放平台"]),
        ("用户操作全旅程", ["注册→创建→配置→测试→发布→运营→迭代", "关键情绪转折：空白画布(intimidating)→模板选择(释然)", "模拟测试消除不知对不对焦虑", "版本管理消除改了变差恐惧"]),
        ("技术架构全景", ["前端：React Flow画布+Zustand状态管理+WebRTC(V2.0)", "引擎：Node.js自研流程解释器+Redis缓存+锁", "通信：Twilio/SIP网关/WebRTC网关三种接入", "AI：LLM网关(多模型)+Deepgram ASR+ElevenLabs TTS+Pinecone RAG", "数据：PostgreSQL(JSONB)+Redis+S3录音+Pinecone向量"]),
        ("功能优先级矩阵", ["Q1高价值低复杂度(P0)：核心层8节点+逻辑层4节点+画布+文字模拟+基础看板", "Q2高价值高复杂度(P1)：AI3节点+集成4节点+真实通话模拟+版本管理+热力图", "Q3低价值低复杂度(P2)：并行分支+自定义看板", "Q4低价值高复杂度(P2)：情感检测+智能路由+多Agent"]),
        ("核心收获", ["AI不是越多越好——三个质变点就够了，规则优先AI增强是正确策略", "文字模拟测试是被低估的杀手功能，秒级迭代解决10分钟可用命题", "节点级运营洞察比AI更有商业价值——从工具到平台的关键升级", "WebRTC开辟无电话号码新市场", "版本管理不是工程功能而是业务保障——消除迭代恐惧"]),
        ("下一步行动", ["MVP开发启动：核心层+逻辑层+画布+平台托管号码+文字模拟", "种子用户招募：50家中小企业内测", "AI节点V1.0验证：意图检测+知识库问答+通话摘要", "通信接入扩展：自带号码(SIP)+真实通话模拟"]),
    ]

    for i, (title, bullets, *rest) in enumerate(slides_data):
        subtitle = rest[0] if rest else None
        if i == 0:
            # Cover slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_slide_bg(slide)
            add_accent_line(slide, 2.5, 2.8, 5.0)
            add_text_box(slide, 1.0, 2.0, 8.0, 1.0, title, font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
            add_text_box(slide, 1.0, 3.0, 8.0, 0.6, subtitle or "", font_size=18, color=LIGHT_TEXT, alignment=PP_ALIGN.CENTER)
            add_text_box(slide, 1.0, 4.5, 8.0, 0.5, "2026年5月", font_size=14, color=LIGHT_TEXT, alignment=PP_ALIGN.CENTER)
        else:
            create_slide(prs, title, bullets, subtitle)

    out_path = os.path.join(os.path.dirname(__file__), "VoiceFlow-Builder-CN.pptx")
    prs.save(out_path)
    print(f"Chinese PPT saved to: {out_path}")

if __name__ == "__main__":
    main()
