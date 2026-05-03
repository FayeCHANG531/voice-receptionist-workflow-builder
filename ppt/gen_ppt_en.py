#!/usr/bin/env python3
"""Generate VoiceFlow Builder PPT (English version)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

BG = RGBColor(0x1e, 0x29, 0x3b)
WHITE = RGBColor(0xff, 0xff, 0xff)
ACCENT = RGBColor(0x4f, 0x46, 0xe5)
LIGHT_TEXT = RGBColor(0x94, 0xa3, 0xb8)

def set_slide_bg(slide, color=BG):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=16, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size); p.font.bold = bold; p.font.color.rgb = color; p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(font_size); p.font.color.rgb = color; p.space_after = Pt(8)
    return txBox

def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.04))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

def create_slide(prs, title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide)
    add_accent_line(slide, 0.8, 1.0, 1.5)
    add_text_box(slide, 0.8, 1.1, 8.5, 0.8, title, font_size=28, bold=True, color=WHITE)
    if subtitle: add_text_box(slide, 0.8, 1.8, 8.5, 0.5, subtitle, font_size=14, color=LIGHT_TEXT)
    add_bullet_list(slide, 1.0, 2.5 if subtitle else 2.2, 8.3, 4.5, bullets, font_size=15, color=WHITE)

def main():
    prs = Presentation(); prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)
    slides = [
        ("VoiceFlow Builder", ["Build a Professional Voice Receptionist in 10 Minutes","Visual Workflow Builder for SMBs","May 2026"], "Visual Workflow Builder SaaS for SMBs"),
        ("The Voice Reception Dilemma for SMBs", ["78% of SMBs still use manual call answering","Existing solutions: too expensive or too basic","Non-technical users can't use developer tools","Core need: usable in 10 min + scalable for complex business"]),
        ("Three Core Capabilities of VoiceFlow", ["Visual workflow builder: drag-and-drop nodes + connections, zero code","Rules-first, AI-enhanced: core flow at zero latency/cost, AI at 3 inflection points","Node-level operational insights: heatmap + funnel for churn attribution"]),
        ("Who Needs VoiceFlow?", ["Small business owners (1-20): no tech team, need quick setup","Mid-size team managers (20-200): need call routing and automation","Online business operators: need website-embedded voice calls"]),
        ("How Does the Workflow Run?", ["Call trigger → Node execution → Conditional routing → Action → End","Each node: Input → Processing → Output, deterministic execution","AI nodes have fallback: auto-revert to rules when AI unavailable"]),
        ("Build a Voice Receptionist in 3 Steps", ["1. Drag nodes: 19 node types from palette, covering all scenarios","2. Connect & configure: connections express data flow, double-click to configure","3. Simulate & publish: text simulation → one-click publish"]),
        ("19 Nodes, Four Progressive Layers", ["Core (8): Call trigger/Greeting/Ask/Collect/Transfer/SMS/Log/End","Logic (4): Conditional/If-Else/Business Hours/Call Type Detection","Integration (4): CRM Lookup/Calendar Booking/Notification/Webhook","AI Enhanced (3): Intent Detection/Knowledge Q&A/Call Summary"]),
        ("Core Layer — The Call Flow Backbone", ["Call Trigger: single entry point, auto-extract caller info","Greeting: time-based script selection, variable substitution","Ask Question: dual input (speech + DTMF), configurable timeout/retry","Transfer Call: warm/cold transfer, auto-degrade on failure","Every node has explicit config items"]),
        ("Logic Layer — Flow Control Hub", ["Conditional Branch: multi-condition routing (≤4 branches), regex support","If/Else: simple binary judgment, reduced config complexity","Business Hours: timezone + holiday calendar + special date override","Call Type Detection: returning/VIP/blacklist auto-identification"]),
        ("Rules-First, AI-Enhanced Strategy", ["17 rule nodes cover 80% scenarios at zero latency/cost","AI only at 3 inflection points: intent detection(-35% call time)/knowledge Q&A(-68% transfers)/call summary(-100% manual)","Emotion detection & smart routing in P2: rule approximation manageable (~15pp)","AI nodes must have fallback: unrecognized intent → DTMF menu"]),
        ("Three AI Inflection Points", ["AI Intent Detection: replace DTMF with NLU, 1 step vs 4.2 steps, 6% vs 28% abandonment","Knowledge Q&A (RAG): one KB covers hundreds of FAQs, blocks 68% human transfers","Call Summary: offline generation, zero real-time latency, 3 min → 30 sec"]),
        ("Cost-Latency Tradeoff Boundaries", ["Pure rules engine: $0 incremental, <10ms latency","+AI Intent Detection: +$0.02/min, +300ms, call time -35%","+Knowledge Q&A: +$0.03/min, +600ms, human transfer -68%","Full AI Agent: +$0.08-0.12/min, +800ms, only for high-complexity"]),
        ("Triple-Mode Telephony Access", ["Platform Managed: out-of-box, purchase→bind→launch, $1-3/mo + $0.01/min","Bring Your Own (SIP/Port): keep existing number, SIP or LNP","WebRTC Browser Call: no phone number needed, embed in website, $0.003-0.008/min","Only platform offering all three modes simultaneously"]),
        ("Two Simulation Modes", ["Text Simulation: instant validation in canvas, second-level iteration","Live Call Test: temp number voice test, real TTS+ASR validation","Underestimated killer feature: reduces design-test loop from minutes to seconds"]),
        ("Safe Iteration, Zero Fear", ["Auto-snapshot on every save, keep latest 50 versions","Version Diff: added green/deleted red/modified yellow","One-click rollback: recover in 30 seconds","Draft/Published separation: published versions cannot be directly modified"]),
        ("Node-Level Insights, Precise Attribution", ["Basic metrics: call volume/answer rate/avg duration/satisfaction","Node heatmap: green>90%/yellow 70-90%/red<70%, locate churn nodes","Conversion funnel: inbound→intent→collection→transfer→end","Exception alerts: consecutive failures/timeout/routing failure real-time"]),
        ("Example: Standard Customer Service Routing", ["Call→Business Hours→(Open) AI Intent→Consult/Book/Complain","Non-business→Knowledge FAQ or Voicemail","Booking→Calendar Booking→SMS Confirmation→End","Key config: 3 intents + 0.7 threshold + DTMF fallback"]),
        ("Example: Medical Clinic Advanced Automation", ["Call→Call Type→(Returning) CRM Lookup→AI Intent→Booking/Consult/Emergency","Emergency→warm transfer to ER, skip all flow","Booking→match doctor's calendar→conflict suggest nearest","Knowledge base covers FAQ (hours/insurance/parking/departments)"]),
        ("Product Positioning & Competitiveness", ["Differentiator 1: Triple-mode telephony + dual-mode simulation (market-unique)","Differentiator 2: Rules-first AI-enhanced layered node system","Differentiator 3: Node-level operational insights (heatmap + funnel)","vs Bland/Vapi/Retell: non-technical users productive in 10 minutes","vs Synthflow: transparent node system + triple-mode + node-level insights"]),
        ("Product Blueprint & Roadmap", ["MVP: Core+Logic + Canvas + Platform numbers + Text sim + Basic dashboard","V1.0: +AI 3 nodes + Integration 4 + BYO numbers + Live test + Version mgmt + Heatmap","V2.0: +WebRTC + Emotion detection + Smart routing + Multi-tenant RBAC + API platform"]),
        ("End-to-End User Journey", ["Signup→Create→Configure→Test→Publish→Operate→Iterate","Key emotion shift: blank canvas (intimidating) → template (relief)","Simulation eliminates 'not sure if it works' anxiety","Version management eliminates 'changes made it worse' fear"]),
        ("Technology Architecture Overview", ["Frontend: React Flow + Zustand + WebRTC (V2.0)","Engine: Node.js flow interpreter + Redis cache/lock","Telephony: Twilio / SIP gateway / WebRTC gateway","AI: LLM gateway + Deepgram ASR + ElevenLabs TTS + Pinecone RAG","Data: PostgreSQL (JSONB) + Redis + S3 + Pinecone vectors"]),
        ("Feature Priority Matrix", ["Q1 High value Low complexity (P0): Core 8 + Logic 4 + Canvas + Text sim + Dashboard","Q2 High value High complexity (P1): AI 3 + Integration 4 + Live test + Version + Heatmap","Q3 Low value Low complexity (P2): Parallel branch + Custom dashboard","Q4 Low value High complexity (P2): Emotion detection + Smart routing + Multi-agent"]),
        ("Key Takeaways", ["Less AI is more — 3 inflection points suffice, rules-first AI-enhanced is right","Text simulation is underestimated killer feature, second-level iteration","Node-level insights > AI commercially — key upgrade from tool to platform","WebRTC opens 'no phone number' new market","Version management is business safeguard, not engineering feature"]),
        ("Next Steps", ["MVP development kickoff: Core+Logic + Canvas + Platform numbers + Text simulation","Seed user recruitment: 50 SMBs for closed beta","AI node V1.0 validation: Intent detection + Knowledge Q&A + Call summary","Telephony expansion: BYO numbers (SIP) + Live call simulation"]),
    ]
    for i, (title, bullets, *rest) in enumerate(slides):
        subtitle = rest[0] if rest else None
        if i == 0:
            slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide)
            add_accent_line(slide, 2.5, 2.8, 5.0)
            add_text_box(slide, 1.0, 2.0, 8.0, 1.0, title, font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
            add_text_box(slide, 1.0, 3.0, 8.0, 0.6, subtitle or "", font_size=18, color=LIGHT_TEXT, alignment=PP_ALIGN.CENTER)
            add_text_box(slide, 1.0, 4.5, 8.0, 0.5, "May 2026", font_size=14, color=LIGHT_TEXT, alignment=PP_ALIGN.CENTER)
        else:
            create_slide(prs, title, bullets, subtitle)

    out_path = os.path.join(os.path.dirname(__file__), "VoiceFlow-Builder-EN.pptx")
    prs.save(out_path)
    print(f"English PPT saved to: {out_path}")

if __name__ == "__main__":
    main()
